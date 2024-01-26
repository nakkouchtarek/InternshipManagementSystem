from bson.binary import Binary
from docx import Document
from langdetect import detect
from pymongo import MongoClient
import spacy
from deep_translator import GoogleTranslator
import pdfplumber
from dateutil import parser
from pptx import Presentation
from TopicHandler import TopicHandler

class FileProcessor:
    def __init__(self):
        self.handler = TopicHandler()
        self.collection = MongoClient("mongodb://localhost:27017").ftest.Internships

    def read_powerpoint(self, file_path):
        presentation = Presentation(file_path)

        text_content = []

        for slide in presentation.slides:
            slide_text = ""

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

            text_content.append(slide_text)

        return text_content

    def detect_encoding_and_language(self, file_path):
        ext = file_path.split(".")[-1]

        if ext == "pdf":
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    language = detect(page.extract_text())
                    return language
        
        elif ext == "docx":
            doc = Document(file_path)
            return detect(doc.paragraphs[1].text)
        
        elif ext == "pptx":
            presentation = Presentation(file_path)

            for slide in presentation.slides:
                slide_text = ""

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"

                return detect(slide_text)

    def file_extract(self, file_path, studentId, reportId):
        paragraphs = []
        ext = file_path.split(".")[-1]
        language = self.detect_encoding_and_language(file_path)

        # read file
        if ext == "pdf":

            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    paragraphs.append(page.extract_text())

        elif ext == "docx":

            doc = Document(file_path)
            paragraphs = [paragraph.text for paragraph in doc.paragraphs]

        elif ext == "pptx":

            paragraphs = self.read_powerpoint(file_path)

        paragraphs = paragraphs[:10]
        
        if language != 'en': paragraphs = [self.translate_phrase(paragraph) for paragraph in paragraphs]

        topic = ''
        topic = self.handler.predict_topic(' '.join(paragraphs))
        
        nlp = spacy.load('en_core_web_lg')

        intern_name = enterprise_name = ''
        duration = ('','')
        supervisors = set()
        companies = {}

        for paragraph in paragraphs:
            lines = paragraph.split("\n")

            concatenated_lines = []

            # grab each line that ends with ":" and concatenates it with the next line
            for i in range(len(lines) - 1):
                if lines[i].endswith(":"):
                    concatenated_line = lines[i] + lines[i + 1]
                    concatenated_lines.append(concatenated_line)
                else:
                    concatenated_lines.append(lines[i])

            concatenated_lines.append(lines[-1])

            lines = concatenated_lines

            for line in lines:
                
                doc_nlp = nlp(line)

                if "supervis" in line.lower():
                    
                    for ent in doc_nlp.ents:
                        if ent.label_ == "PERSON":
                            supervisors.add(ent.text)

                if "directed" in line.lower() or "presented" in line.lower():
                    
                    if intern_name == '':

                        for ent in doc_nlp.ents:

                            if ent.label_ == "PERSON":
                                intern_name = ent.text
                                break

                for ent in doc_nlp.ents:
                    
                    if ent.label_ == "ORG":

                        company = ent.text.strip()

                        if company not in companies:
                            companies[company] = 1
                        else:
                            companies[company] += 1

                    elif ent.label_ == "DATE":
                        
                        if duration == ('','') and not ent.text.startswith("to") and "to" in ent.text.lower():
                            date = ent.text.strip().split('to')

                            for i, d in enumerate(date):
                                date[i] = [elem.replace(',','') for elem in d.split(" ") if elem != '']
                            
                            if len(date[0]) == 3:
                                date[1][2].append(date[0][2])
                            elif len(date[1]) == 3:
                                date[0].append(date[1][2])

                            date = [' '.join(d) for d in date]

                            start_date = parser.parse(date[0])
                            end_date = parser.parse(date[1])

                            duration = (str(start_date).split(" ")[0], str(end_date).split(" ")[0])

        enterprise_name = sorted(companies.items(),reverse=True, key=lambda x:x[1])[0][0]
        supervisors = [name for name in supervisors if name.lower().strip() != intern_name.lower().strip()]

        self.collection.insert_one({
            "student_id":studentId,
            "company_name":enterprise_name,
            "start_date":duration[0],
            "end_date":duration[1],
            "supervisor_name":' ,'.join(supervisors),
            "theme": topic,
            "reportId": reportId
        })

    def translate_phrase(self, phrase):
        return GoogleTranslator(source='auto', target='en').translate(phrase)